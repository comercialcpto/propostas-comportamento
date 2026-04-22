import streamlit as st
import pandas as pd
import math
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import io
import datetime

# ==========================================
# 1. FUNÇÕES AUXILIARES E MATEMÁTICA
# ==========================================
VERDE_CPTO = RGBColor(0, 153, 116) 
CINZA_ESCURO = RGBColor(64, 64, 64)

def formatar_moeda(valor):
    return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")

def formatar_celula_tabela(cell, texto):
    cell.text = str(texto)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "DIN Alternate"
            run.font.size = Pt(14)

def valor_por_extenso(valor):
    unidades = ["", "um", "dois", "três", "quatro", "cinco", "seis", "sete", "oito", "nove"]
    dez_a_dezenove = ["dez", "onze", "doze", "treze", "quatorze", "quinze", "dezesseis", "dezessete", "dezoito", "dezenove"]
    dezenas = ["", "", "vinte", "trinta", "quarenta", "cinquenta", "sessenta", "setenta", "oitenta", "noventa"]
    centenas = ["", "cento", "duzentos", "trezentos", "quatrocentos", "quinhentos", "seiscentos", "setecentos", "oitocentos", "novecentos"]

    def converter_trio(n):
        if n == 100: return "cem"
        if n == 0: return ""
        c, resto = n // 100, n % 100
        d, u = resto // 10, resto % 10
        partes = []
        if c > 0: partes.append(centenas[c])
        if d == 1: partes.append(dez_a_dezenove[u])
        else:
            if d > 1: partes.append(dezenas[d])
            if u > 0: partes.append(unidades[u])
        return " e ".join(partes)

    inteiro = int(valor)
    if inteiro == 0: return "zero reais"
    milhoes = inteiro // 1000000
    milhares = (inteiro % 1000000) // 1000
    resto = inteiro % 1000
    
    resultado = []
    if milhoes > 0: resultado.append(converter_trio(milhoes) + (" milhões" if milhoes > 1 else " milhão"))
    if milhares > 0: resultado.append(converter_trio(milhares) + " mil")
    if resto > 0: resultado.append(converter_trio(resto))
        
    extenso_final = ", ".join(resultado).replace(", e", " e")
    return extenso_final.capitalize() + (" reais" if inteiro > 1 else " real")

# --- MOTOR ESTATÍSTICO DCS ---
def calcular_amostra(N, margem_erro, proporcao, z=1.96):
    if N <= 0: return 0
    numerador = N * (z**2) * proporcao * (1 - proporcao)
    denominador = (margem_erro**2) * (N - 1) + (z**2) * proporcao * (1 - proporcao)
    return math.ceil(numerador / denominador)

def calcular_amortizacao(qtd_parcelas):
    pesos_base = [20, 20, 10, 10] + [5] * 30
    pesos_ativos = pesos_base[:qtd_parcelas]
    soma = sum(pesos_ativos)
    percentuais = [round((p / soma) * 100) for p in pesos_ativos]
    diferenca = 100 - sum(percentuais)
    if diferenca != 0: percentuais[0] += diferenca
    return percentuais

# ==========================================
# 2. MOTOR DE PROCESSAMENTO PPTX
# ==========================================
def deletar_slide(prs, slide):
    id_dict = { s.id: [i, s.rId] for i, s in enumerate(prs.slides._sldIdLst) }
    prs.part.drop_rel(id_dict[slide.slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide.slide_id][0]]

def remover_linha_tabela(table, row_idx):
    try:
        tr = table.rows[row_idx]._tr
        tr.getparent().remove(tr)
    except Exception: pass

def remover_coluna_tabela(table, col_idx):
    try:
        tbl = table._tbl
        grid = tbl.tblGrid
        col = grid.gridCol_lst[col_idx]
        grid.remove(col)
        for tr in tbl.tr_lst:
            tc = tr.tc_lst[col_idx]
            tr.remove(tc)
    except Exception: pass

def processar_apresentacao(template_file, mapa, atividades, tipo_doc, dados_fin=None, qtd_meses=12):
    prs = Presentation(template_file)
    slides_para_deletar = []

    for slide in prs.slides:
        deletar_este_slide = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and "Para DCS" in shape.text:
                if mapa.get("{{SERVICO}}", "") != "Diagnóstico (DCS/Clima/DCMA)":
                    deletar_este_slide = True
                    break 

            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in mapa.items():
                            if key in run.text: run.text = run.text.replace(key, str(value))
            
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for p in cell.text_frame.paragraphs:
                                for run in p.runs:
                                    for key, value in mapa.items():
                                        if key in run.text: run.text = run.text.replace(key, str(value))

                if len(tbl.columns) >= 12 and len(atividades) > 0:
                    colunas_para_deletar = list(range(qtd_meses + 1, len(tbl.columns)))
                    try:
                        largura_original = shape.width
                        largura_deletada = sum([tbl.columns[c].width for c in colunas_para_deletar])
                        nova_largura = largura_original - largura_deletada
                        shape.left = int((prs.slide_width - nova_largura) / 2)
                    except Exception: pass

                    for c_idx in reversed(colunas_para_deletar):
                        remover_coluna_tabela(tbl, c_idx)

                    linhas_deletar = list(range(len(atividades) + 1, len(tbl.rows)))
                    for r_idx in reversed(linhas_deletar):
                        remover_linha_tabela(tbl, r_idx)

                    for row_idx, atividade in enumerate(atividades):
                        target_row = row_idx + 1 
                        if target_row < len(tbl.rows):
                            row = tbl.rows[target_row]
                            cell = row.cells[0]
                            cell.text = atividade['nome']
                            
                            tamanho_str = len(atividade['nome'])
                            fonte_tamanho = 12
                            if tamanho_str > 60: fonte_tamanho = 8
                            elif tamanho_str > 40: fonte_tamanho = 9
                            elif tamanho_str > 20: fonte_tamanho = 10

                            if cell.text_frame.paragraphs:
                                p = cell.text_frame.paragraphs[0]
                                if p.runs:
                                    run = p.runs[0]
                                    run.font.name = "Calibri"
                                    run.font.size = Pt(fonte_tamanho)
                                    run.font.color.rgb = CINZA_ESCURO

                            for m_idx in range(1, len(tbl.columns)):
                                if m_idx in atividade['meses']:
                                    cell_mes = row.cells[m_idx]
                                    cell_mes.fill.solid()
                                    cell_mes.fill.fore_color.rgb = VERDE_CPTO

                if tipo_doc == "Comercial" and dados_fin:
                    try:
                        cabecalho = tbl.rows[0].cells[0].text.strip().lower()
                        
                        if "macro" in cabecalho:
                            acoes = dados_fin['acoes']
                            linhas_para_deletar = []
                            for idx in range(1, len(tbl.rows)):
                                cell_text = tbl.rows[idx].cells[0].text.strip().lower()
                                if "investimento total" in cell_text:
                                    formatar_celula_tabela(tbl.rows[idx].cells[1], formatar_moeda(dados_fin['total_op1']))
                                    formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(dados_fin['total_op2']))
                                elif idx <= len(acoes):
                                    formatar_celula_tabela(tbl.rows[idx].cells[0], acoes[idx-1]['nome'])
                                    formatar_celula_tabela(tbl.rows[idx].cells[1], formatar_moeda(acoes[idx-1]['v1']))
                                    formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(acoes[idx-1]['v2']))
                                else:
                                    linhas_para_deletar.append(idx)
                            for idx in reversed(linhas_para_deletar):
                                remover_linha_tabela(tbl, idx)

                        elif "meses" in cabecalho:
                            parcelas = dados_fin['parcelas']
                            linhas_para_deletar = []
                            for idx in range(1, len(tbl.rows)):
                                cell_text = tbl.rows[idx].cells[0].text.strip().lower()
                                if "total" in cell_text and "investimento" not in cell_text:
                                    formatar_celula_tabela(tbl.rows[idx].cells[1], "100%")
                                    formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(dados_fin['total_op2']))
                                elif idx <= len(parcelas):
                                    formatar_celula_tabela(tbl.rows[idx].cells[0], f"M{idx}")
                                    formatar_celula_tabela(tbl.rows[idx].cells[1], f"{parcelas[idx-1]}%")
                                    val_calc = dados_fin['total_op2'] * (parcelas[idx-1] / 100)
                                    formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(val_calc))
                                else:
                                    linhas_para_deletar.append(idx)
                            for idx in reversed(linhas_para_deletar):
                                remover_linha_tabela(tbl, idx)
                    except Exception: pass 

        if deletar_este_slide: slides_para_deletar.append(slide)

    for slide in slides_para_deletar: deletar_slide(prs, slide)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# ==========================================
# 3. INTERFACE DE USUÁRIO E ESTADOS
# ==========================================
st.set_page_config(page_title="Sistema Comportamento", layout="wide")

if "pptx_gerado" not in st.session_state: st.session_state.pptx_gerado = None
if "nome_arquivo" not in st.session_state: st.session_state.nome_arquivo = ""
if 'tentou_gerar' not in st.session_state: st.session_state.tentou_gerar = False

if 'unidades_dcs' not in st.session_state:
    st.session_state.unidades_dcs = [{"id": 0, "nome": "Unidade Matriz", "pop_total": 0, "lideres": 0}]

# --- NOVO ESTADO: FASES DO PLANO DETALHADO ---
if 'fases_dcs' not in st.session_state:
    st.session_state.fases_dcs = [
        {"id": 0, "nome": "Planejamento / Reunião de Abertura", "horas": 4},
        {"id": 1, "nome": "Workshop - Gestão de Cultura de Segurança", "horas": 2},
        {"id": 2, "nome": "Elaboração do Relatório", "horas": 56},
        {"id": 3, "nome": "Apresentação dos Resultados", "horas": 4},
        {"id": 4, "nome": "Plano de Transformação Cultural (PTC)", "horas": 12},
        {"id": 5, "nome": "Suporte e Acompanhamento", "horas": 64}
    ]

def acionar_geracao():
    st.session_state.tentou_gerar = True

def adicionar_unidade():
    novo_id = len(st.session_state.unidades_dcs)
    st.session_state.unidades_dcs.append({"id": novo_id, "nome": f"Nova Unidade {novo_id+1}", "pop_total": 0, "lideres": 0})

def remover_unidade(id_remover):
    if len(st.session_state.unidades_dcs) > 1:
        st.session_state.unidades_dcs = [u for u in st.session_state.unidades_dcs if u['id'] != id_remover]

def adicionar_fase():
    novo_id = max([f['id'] for f in st.session_state.fases_dcs], default=-1) + 1
    st.session_state.fases_dcs.append({"id": novo_id, "nome": "Nova Etapa", "horas": 0})

def remover_fase(id_remover):
    st.session_state.fases_dcs = [f for f in st.session_state.fases_dcs if f['id'] != id_remover]


st.sidebar.title("🧭 Navegação")
menu = st.sidebar.radio("Selecione o Módulo:", ["🏠 Início", "💰 Criar Precificação", "📊 Criar Apresentação"])

# ==========================================
# MÓDULO 1: INÍCIO
# ==========================================
if menu == "🏠 Início":
    st.title("Bem-vindo ao Sistema Comercial - Comportamento")
    st.write("Selecione um módulo no menu lateral para começar.")
    st.session_state.tentou_gerar = False

# ==========================================
# MÓDULO 2: PRECIFICAÇÃO
# ==========================================
elif menu == "💰 Criar Precificação":
    st.title("💰 Motor de Precificação")
    servico_prec = st.selectbox("Selecione o Serviço para Precificar:", ["Diagnóstico (DCS)", "Mapeamento de Liderança (MPL) - Em breve"])
    
    if servico_prec == "Diagnóstico (DCS)":
        
        # --- 1. CADASTRO DE UNIDADES ---
        st.markdown("### 📊 1. Cadastro de População e Unidades")
        st.info("Insira as áreas ou unidades. O sistema fará o cálculo da amostra usando 95% de confiança.")
        
        for i, und in enumerate(st.session_state.unidades_dcs):
            c1, c2, c3, c4 = st.columns([0.4, 0.2, 0.2, 0.2])
            und['nome'] = c1.text_input(f"Nome da Unidade/Área", value=und['nome'], key=f"nome_{und['id']}")
            und['pop_total'] = c2.number_input("População Total", min_value=0, value=und['pop_total'], key=f"pop_{und['id']}")
            und['lideres'] = c3.number_input("Total de Líderes", min_value=0, value=und['lideres'], key=f"lid_{und['id']}")
            
            if len(st.session_state.unidades_dcs) > 1:
                if c4.button("🗑️ Remover", key=f"rem_{und['id']}"):
                    remover_unidade(und['id'])
                    st.rerun()

        st.button("➕ Adicionar Unidade / Área", on_click=adicionar_unidade)
        
        st.markdown("---")
        
        # --- 2. MOTOR ESTATÍSTICO ---
        st.markdown("### ⚙️ 2. Motor Estatístico (Coleta de Dados em Campo)")
        
        total_horas_campo = 0
        total_amostra = 0
        
        for und in st.session_state.unidades_dcs:
            if und['pop_total'] > 0:
                amostra_total = calcular_amostra(und['pop_total'], margem_erro=0.08, proporcao=0.5)
                amostra_lideres = calcular_amostra(und['lideres'], margem_erro=0.05, proporcao=0.8)
                
                turmas_hm = math.ceil(amostra_lideres / 12)
                turmas_foco = math.ceil((amostra_total - amostra_lideres) / 12) if amostra_total > amostra_lideres else 0
                horas_hm = turmas_hm * 2
                horas_foco = turmas_foco * 2
                
                entrevistas = max(0, amostra_lideres - (turmas_hm * 12)) if amostra_lideres > 0 else 8
                horas_entrevistas = entrevistas * 1.5
                
                total_horas_unidade = horas_hm + horas_foco + horas_entrevistas
                total_horas_campo += total_horas_unidade
                total_amostra += amostra_total
                
                with st.expander(f"📌 Resultados: {und['nome']} (Amostra: {amostra_total} pessoas)", expanded=True):
                    rc1, rc2, rc3 = st.columns(3)
                    rc1.metric("Amostra Total (Het. 8%)", amostra_total, f"Pop: {und['pop_total']}")
                    rc2.metric("Amostra Líderes (Hom. 5%)", amostra_lideres, f"Pop: {und['lideres']}")
                    rc3.metric("Horas Estimadas de Coleta", f"{total_horas_unidade}h")
                    st.write(f"**Sugestão de Escopo:** {turmas_hm} Turmas H&M | {turmas_foco} Grupos Focais | {entrevistas} Entrevistas Ind.")

        st.markdown("---")

        # --- 3. PLANO DETALHADO (NOVO) ---
        st.markdown("### 📋 3. Plano Detalhado (Etapas Adicionais e Backoffice)")
        st.info("Estas são as etapas padrões do DCS. Modifique, adicione Oficinas ou remova o que não for usar.")

        taxa_hora = st.number_input("Valor da Taxa Hora Técnica (R$)", min_value=0.0, value=480.0, step=10.0, key="taxa_hora_topo")

        total_horas_fases = 0

        for fase in st.session_state.fases_dcs:
            f1, f2, f3 = st.columns([0.6, 0.2, 0.2])
            fase['nome'] = f1.text_input("Nome da Etapa", value=fase['nome'], key=f"fnome_{fase['id']}")
            fase['horas'] = f2.number_input("Carga Horária (h)", min_value=0, value=fase['horas'], key=f"fhoras_{fase['id']}")
            
            total_horas_fases += fase['horas']

            if f3.button("🗑️ Remover", key=f"frem_{fase['id']}"):
                remover_fase(fase['id'])
                st.rerun()

        st.button("➕ Adicionar Etapa no Plano", on_click=adicionar_fase)

        # Resumo do Plano Detalhado
        valor_parcial_fases = total_horas_fases * taxa_hora
        st.markdown(f"**Resumo do Plano Detalhado:** {len(st.session_state.fases_dcs)} etapas cadastradas | **{total_horas_fases} horas** totais.")
        st.caption(f"Equivalente a: {total_horas_fases}h * {formatar_moeda(taxa_hora)} = **{formatar_moeda(valor_parcial_fases)}**")

        st.markdown("---")
        
        # --- 4. PRECIFICACAO BASE E LOGISTICA ---
        st.markdown("### 💰 4. Precificação Final e Logística")
        
        # A soma total de horas agora inclui as horas variáveis de campo + as horas dinâmicas do Plano Detalhado
        horas_totais = total_horas_campo + total_horas_fases
        valor_op1 = horas_totais * taxa_hora
        
        st.write("#### Planejamento de Logística (Opção 2)")
        tipo_logistica = st.selectbox("Formato de cálculo de deslocamento:", [
            "1. Sem Logística (100% pelo Cliente)", 
            "2. Logística Base (Alimentação + Táxi da Base)", 
            "3. Logística Completa (Cotações Detalhadas)", 
            "4. Logística Estimada (Percentual %)"
        ])
        
        logistica_total = 0.0
        
        if tipo_logistica == "1. Sem Logística (100% pelo Cliente)":
            st.info("A Opção 2 terá o mesmo valor da Opção 1, pois o cliente arcará com todos os custos de viagem diretamente.")
            logistica_total = 0.0

        elif tipo_logistica == "2. Logística Base (Alimentação + Táxi da Base)":
            c_ida, c_dia = st.columns(2)
            qtd_idas = c_ida.number_input("Quantidade de Idas Presenciais", min_value=1, value=1)
            dias_ida = c_dia.number_input("Dias por Ida (padrão 5 = semana cheia)", min_value=1, value=5)
            
            custo_taxi = qtd_idas * (150 * 2)
            custo_alimentacao = qtd_idas * dias_ida * 120
            logistica_total = custo_taxi + custo_alimentacao
            
            st.success(f"**Cálculo Base:** {qtd_idas} idas (Táxi: R$ {custo_taxi:,.2f}) + {qtd_idas * dias_ida} dias de alimentação (R$ {custo_alimentacao:,.2f}) = **R$ {logistica_total:,.2f}**")

        elif tipo_logistica == "3. Logística Completa (Cotações Detalhadas)":
            c_ida, c_dia = st.columns(2)
            qtd_idas = c_ida.number_input("Quantidade de Idas Presenciais", min_value=1, value=1)
            dias_ida = c_dia.number_input("Dias por Ida (padrão 5)", min_value=1, value=5)
            
            st.markdown("##### 🏨 Hospedagem")
            ch1, ch2 = st.columns(2)
            hotel_barato = ch1.number_input("Diária - Hotel Mais Barato (R$)", min_value=0.0, step=10.0)
            hotel_caro = ch2.number_input("Diária - Hotel Mais Caro (R$)", min_value=0.0, step=10.0)
            media_hotel = (hotel_barato + hotel_caro) / 2
            custo_hotel = media_hotel * dias_ida * qtd_idas
            
            st.markdown("##### ✈️ Passagens Aéreas")
            st.caption("A fórmula = ((Barato + Caro) * 3.2) / 6 será aplicada automaticamente.")
            qtd_trechos = st.number_input("Quantos trechos aéreos diferentes?", min_value=0, value=1)
            custo_aereo_total = 0.0
            for i in range(int(qtd_trechos)):
                ca1, ca2 = st.columns(2)
                aereo_barato = ca1.number_input(f"Trecho {i+1} - Mais Barato (R$)", min_value=0.0, step=50.0, key=f"ab_{i}")
                aereo_caro = ca2.number_input(f"Trecho {i+1} - Mais Caro (R$)", min_value=0.0, step=50.0, key=f"ac_{i}")
                media_aereo = ((aereo_barato + aereo_caro) * 3.2) / 6
                custo_aereo_total += (media_aereo * qtd_idas)

            st.markdown("##### 🚗 Locação de Veículo e Deslocamentos")
            cv1, cv2 = st.columns(2)
            diaria_carro = cv1.number_input("Valor da Diária do Carro (R$)", min_value=0.0, step=10.0)
            custo_carro = diaria_carro * dias_ida * qtd_idas
            
            dist_hotel_cliente = cv2.number_input("Distância Hotel ⇄ Cliente (Km - Total Ida e Volta diária)", min_value=0.0, step=5.0)
            custo_combustivel_hotel = (dist_hotel_cliente / 9.0) * 6.0 * dias_ida * qtd_idas
            
            st.markdown("##### 🛣️ Aeroporto até o Cliente")
            cae1, cae2 = st.columns(2)
            dist_aero_cliente = cae1.number_input("Dist. Aeroporto ⇄ Cliente (Km - Total Ida e Volta da viagem)", min_value=0.0, step=10.0)
            pedagio_aero = cae2.number_input("Valor Pedágios Aeroporto ⇄ Cliente (R$ Total)", min_value=0.0, step=5.0)
            
            custo_combustivel_aero = (dist_aero_cliente / 9.0) * 6.0 * qtd_idas
            custo_pedagio = pedagio_aero * qtd_idas
            
            logistica_total = custo_hotel + custo_aereo_total + custo_carro + custo_combustivel_hotel + custo_combustivel_aero + custo_pedagio
            
            st.info(f"**Resumo Logística Cotação:** Hospedagem (R$ {custo_hotel:,.2f}) + Aéreo (R$ {custo_aereo_total:,.2f}) + Carro (R$ {custo_carro:,.2f}) + Combustível (R$ {(custo_combustivel_hotel+custo_combustivel_aero):,.2f}) + Pedágio (R$ {custo_pedagio:,.2f}) = **R$ {logistica_total:,.2f}**")

        elif tipo_logistica == "4. Logística Estimada (Percentual %)":
            perc_logistica = st.number_input("Margem Estimada de Logística (%)", min_value=0, max_value=100, value=30)
            logistica_total = valor_op1 * (perc_logistica / 100)
            st.success(f"Cálculo: {perc_logistica}% sobre o valor técnico da OP1 = **R$ {logistica_total:,.2f}**")
        
        valor_op2 = valor_op1 + logistica_total
        
        st.markdown("---")
        st.success(f"**Carga Horária Total (Amostragem de Campo + Etapas do Plano):** {horas_totais} horas")
        c_tot1, c_tot2 = st.columns(2)
        c_tot1.metric("Total OP1 (Serviço Técnico)", formatar_moeda(valor_op1))
        c_tot2.metric("Total OP2 (Com Logística)", formatar_moeda(valor_op2))

# ==========================================
# MÓDULO 4: CRIAR APRESENTAÇÃO 
# ==========================================
elif menu == "📊 Criar Apresentação":
    st.title("📊 Gerador de Propostas")
    tipo_apresentacao = st.radio("Selecione o tipo de documento:", ["Apresentação Técnica", "Apresentação Comercial"], horizontal=True)
    
    with st.sidebar:
        st.markdown("---")
        st.header("📁 Template Original")
        template_upload = st.file_uploader(f"Suba o template ({tipo_apresentacao})", type="pptx")
    
    campos_vazios = []

    with st.expander("📍 1. Identificação Geral", expanded=True):
        c1, c2, c3 = st.columns(3)
        servico = c1.selectbox("Serviço Principal", ["Diagnóstico (DCS/Clima/DCMA)", "Mapeamento de Liderança (MPL)", "Riscos Psicossociais (RPS)", "Pulse", "Pontuais / Palestras"])
        
        cliente = c2.text_input("Nome da Empresa ({{CLIENTE}})*")
        if st.session_state.tentou_gerar and not cliente: 
            c2.error("Campo obrigatório!")
            campos_vazios.append("Nome da Empresa")

        unidade = c3.text_input("Unidade ({{UNIDADE}})*")
        if st.session_state.tentou_gerar and not unidade: 
            c3.error("Campo obrigatório!")
            campos_vazios.append("Unidade")
        
        c4, c5, c6 = st.columns(3)
        num_prop = c4.text_input("Nº da Proposta ({{NUM_PROP}})*")
        if st.session_state.tentou_gerar and not num_prop: 
            c4.error("Campo obrigatório!")
            campos_vazios.append("Nº da Proposta")

        escopo_tag = c5.text_input("Título do Escopo ({{ESCOPO}})*")
        if st.session_state.tentou_gerar and not escopo_tag: 
            c5.error("Campo obrigatório!")
            campos_vazios.append("Título do Escopo")

        prazo = c6.text_input("Prazo ({{PRAZO}})*")
        if st.session_state.tentou_gerar and not prazo: 
            c6.error("Campo obrigatório!")
            campos_vazios.append("Prazo")
        
        c7, c8, c9 = st.columns(3)
        formato = c7.selectbox("Formato ({{FORMATO}})*", ["Híbrido", "Presencial", "Online"])
        
        idiomas_selecionados = c8.multiselect("Idioma ({{IDIOMA}})*", ["Português", "Espanhol", "Inglês"], default=["Português"])
        if len(idiomas_selecionados) == 1: idioma_str = idiomas_selecionados[0]
        elif len(idiomas_selecionados) == 2: idioma_str = f"{idiomas_selecionados[0]} e {idiomas_selecionados[1]}"
        elif len(idiomas_selecionados) > 2: idioma_str = ", ".join(idiomas_selecionados[:-1]) + f" e {idiomas_selecionados[-1]}"
        else: 
            idioma_str = ""
            if st.session_state.tentou_gerar:
                c8.error("Selecione pelo menos um idioma!")
                campos_vazios.append("Idioma")
        
        idas = c9.number_input("Nº de Idas Presenciais ({{IDAS}})", min_value=0, value=0)
        
        justificativa = st.text_area("Justificativa ({{JUSTIFICATIVA}})*")
        if st.session_state.tentou_gerar and not justificativa: 
            st.error("A Justificativa é obrigatória!")
            campos_vazios.append("Justificativa")

        objetivo = st.text_area("Objetivo ({{OBJETIVO}})*")
        if st.session_state.tentou_gerar and not objetivo: 
            st.error("O Objetivo é obrigatório!")
            campos_vazios.append("Objetivo")

    with st.expander("📅 2. Cronograma de Avanço (Gantt)", expanded=True):
        cg1, cg2 = st.columns(2)
        qtd_fases = cg1.number_input("Quantas Fases?", min_value=1, value=5)
        qtd_meses_projeto = cg2.number_input("Duração total do projeto (meses)", min_value=1, value=12)
        
        atividades_lista = []
        for i in range(qtd_fases):
            ca, cm = st.columns([0.4, 0.6])
            nome_at = ca.text_input(f"Nome da Fase {i+1}", key=f"f_{i}")
            habilitar_meses = len(nome_at) >= 3
            texto_placeholder = "Selecione os meses" if habilitar_meses else "Digite o nome da fase para liberar"
            meses_at = cm.multiselect(texto_placeholder, list(range(1, int(qtd_meses_projeto) + 1)), key=f"m_{i}", disabled=not habilitar_meses)
            
            if habilitar_meses and meses_at:
                atividades_lista.append({"nome": nome_at, "meses": meses_at})

    # ==========================================
    # APRESENTAÇÃO TÉCNICA
    # ==========================================
    if tipo_apresentacao == "Apresentação Técnica":
        with st.expander("👥 3. Detalhamento do Público e Relatórios", expanded=True):
            cp1, cp2, cp3 = st.columns(3)
            n_pr = cp1.number_input("Executivos", min_value=0, value=0)
            n_exec = cp2.number_input("Alta Liderança", min_value=0, value=0)
            n_coord = cp3.number_input("Coordenadores", min_value=0, value=0)
            n_super = cp1.number_input("Supervisores", min_value=0, value=0)
            n_lid_extra = cp2.number_input("Outros Líderes", min_value=0, value=0)
            n_sec = cp3.number_input("Segurança", min_value=0, value=0)
            n_oper = cp1.number_input("Operacional", min_value=0, value=0)
            n_col3 = cp2.number_input("Terceiros", min_value=0, value=0)
            n_lid3 = cp3.number_input("Líderes Terc.", min_value=0, value=0)
            
            n_lid_total = n_pr + n_exec + n_coord + n_super + n_lid_extra
            n_prop = n_lid_total + n_sec + n_oper
            n_p_terc = n_prop + n_col3 + n_lid3

            st.markdown("### 📊 Contador de População (Ao Vivo)")
            m1, m2, m3 = st.columns(3)
            m1.metric("Líderes Totais", n_lid_total)
            m2.metric("Colaboradores Próprios", n_prop)
            m3.metric("Total (Próprios + Terc.)", n_p_terc, delta=f"+{n_col3 + n_lid3} terceiros")

            st.write("---")
            cr1, cr2, cr3 = st.columns(3)
            qtd_rel = cr1.number_input("Qtd de Unidades com Relatório", min_value=0, value=1)
            tem_corp = cr2.checkbox("Relatório Corporativo?")
            tot_rel = qtd_rel + (1 if tem_corp else 0)
            tot_plan = cr3.number_input("Total de PTCs", min_value=0, value=1)

        if st.button("🚀 VALIDAR E GERAR TÉCNICA", on_click=acionar_geracao, key="btn_tecnica"):
            erros = []
            if not template_upload: erros.append("Template Original (Upload)")
            if len(campos_vazios) > 0: erros.extend(campos_vazios)
            if len(atividades_lista) == 0: erros.append("Pelo menos 1 Fase do Cronograma com meses selecionados")

            if erros:
                st.error(f"⚠️ IMPOSSÍVEL GERAR! Preencha os campos obrigatórios: {', '.join(erros)}.")
                st.session_state.pptx_gerado = None
            else:
                mapa = {
                    "{{SERVICO}}": servico, "{{CLIENTE}}": cliente, "{{UNIDADE}}": unidade, 
                    "{{NUM_PROP}}": num_prop, "{{ESCOPO}}": escopo_tag,
                    "{{DATA}}": datetime.date.today().strftime("%d/%m/%Y"),
                    "{{JUSTIFICATIVA}}": justificativa, "{{OBJETIVO}}": objetivo,
                    "{{PUBLICO}}": str(n_p_terc), "{{PRAZO}}": prazo, "{{FORMATO}}": formato, 
                    "{{IDIOMA}}": idioma_str, "{{IDAS}}": str(idas),
                    "{{N_PR}}": str(n_pr), "{{N_EXEC}}": str(n_exec), "{{N_COORD}}": str(n_coord), 
                    "{{N_SUPER}}": str(n_super), "{{N_LID}}": str(n_lid_total), "{{N_SEC}}": str(n_sec), 
                    "{{N_OPER}}": str(n_oper), "{{N_PROP}}": str(n_prop), "{{N_COL3}}": str(n_col3), 
                    "{{N_LID3}}": str(n_lid3), "{{N_PTERC}}": str(n_p_terc),
                    "{{TOT_REL}}": str(tot_rel), "{{QTD_REL}}": str(qtd_rel), "{{TOT_PLAN}}": str(tot_plan)
                }
                
                with st.spinner("Construindo arquivo..."):
                    st.session_state.pptx_gerado = processar_apresentacao(template_upload, mapa, atividades_lista, "Técnica", None, qtd_meses_projeto)
                    st.session_state.nome_arquivo = f"Tecnica_{cliente}.pptx"
                st.success("Técnica gerada com sucesso! Clique abaixo para baixar.")
                st.session_state.tentou_gerar = False

        if st.session_state.pptx_gerado and st.session_state.nome_arquivo.startswith("Tecnica"):
            st.download_button("⬇️ Baixar Documento Gerado", data=st.session_state.pptx_gerado, file_name=st.session_state.nome_arquivo, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    # ==========================================
    # APRESENTAÇÃO COMERCIAL
    # ==========================================
    elif tipo_apresentacao == "Apresentação Comercial":
        with st.expander("👥 3. Público Alvo", expanded=True):
            publico_total = st.number_input("Público Total ({{PUBLICO}})*", min_value=0, value=0)
            if st.session_state.tentou_gerar and publico_total <= 0:
                st.error("Campo obrigatório! Informe o público total.")
                campos_vazios.append("Público Total")

        with st.expander("💰 4. Detalhamento de Investimento e Parcelas", expanded=True):
            modo_logistica = st.radio("Como a Logística será tratada?", ["Estimada (Soma +30% automático)", "Cotada (Informar manualmente)"])
            
            st.write("**Fases de Investimento:**")
            qtd_acoes = st.number_input("Quantas Macro Ações?", min_value=1, value=3)
            
            acoes_fin = []
            total_op1 = 0.0
            total_op2 = 0.0
            
            for i in range(qtd_acoes):
                cf1, cf2, cf3 = st.columns(3)
                n_acao = cf1.text_input(f"Ação {i+1}", key=f"ac_n_{i}")
                v1 = cf2.number_input(f"Valor Opção 1 (R$)", min_value=0.0, key=f"ac_v1_{i}", value=0.0)
                
                if modo_logistica == "Estimada (Soma +30% automático)":
                    v2 = v1 * 1.3
                    ui_v2 = formatar_moeda(v2).replace("$", "\$")
                    cf3.info(f"Opção 2: {ui_v2}")
                else:
                    v2 = cf3.number_input(f"Valor Opção 2 (R$)", min_value=0.0, key=f"ac_v2_{i}", value=0.0)
                
                if n_acao:
                    acoes_fin.append({'nome': n_acao, 'v1': v1, 'v2': v2})
                    total_op1 += v1
                    total_op2 += v2
                elif st.session_state.tentou_gerar:
                    cf1.error("Obrigatório nomear a ação ou remover quantidade.")
                    campos_vazios.append(f"Nome da Ação {i+1}")
            
            st.markdown("---")
            ui_op1 = formatar_moeda(total_op1).replace("$", "\$")
            ui_op2 = formatar_moeda(total_op2).replace("$", "\$")
            st.markdown(f"**Total OP1:** {ui_op1} | **Total OP2:** {ui_op2}")
            
            st.write("**Condições de Pagamento:**")
            qtd_parcelas = st.number_input("Quantidade de Parcelas ({{QTD_PARCELAS}})", min_value=1, value=12)

        if st.button("🚀 VALIDAR E GERAR COMERCIAL", on_click=acionar_geracao, key="btn_comercial"):
            erros = []
            if not template_upload: erros.append("Template Original (Upload)")
            if len(campos_vazios) > 0: erros.extend(campos_vazios)
            if len(atividades_lista) == 0: erros.append("Pelo menos 1 Fase do Cronograma com meses selecionados")
            if len(acoes_fin) == 0: erros.append("Pelo menos 1 Ação Financeira nomeada")

            if erros:
                st.error(f"⚠️ IMPOSSÍVEL GERAR! Preencha os campos obrigatórios: {', '.join(erros)}.")
                st.session_state.pptx_gerado = None
            else:
                mapa_comercial = {
                    "{{SERVICO}}": servico, "{{CLIENTE}}": cliente, "{{UNIDADE}}": unidade, 
                    "{{NUM_PROP}}": num_prop, "{{ESCOPO}}": escopo_tag,
                    "{{DATA}}": datetime.date.today().strftime("%d/%m/%Y"),
                    "{{JUSTIFICATIVA}}": justificativa, "{{OBJETIVO}}": objetivo,
                    "{{PUBLICO}}": str(publico_total), 
                    "{{PRAZO}}": prazo, "{{FORMATO}}": formato, "{{IDIOMA}}": idioma_str, "{{IDAS}}": str(idas),
                    "{{VALOR_OP1}}": formatar_moeda(total_op1),
                    "{{VALOR_OP2}}": formatar_moeda(total_op2),
                    "{{VALOR_OP1_EXT}}": valor_por_extenso(total_op1),
                    "{{VALOR_OP2_EXT}}": valor_por_extenso(total_op2),
                    "{{QTD_PARCELAS}}": str(qtd_parcelas),
                    "{{VLR1_PARCELAS}}": formatar_moeda(total_op1/qtd_parcelas),
                    "{{VLR2_PARCELAS}}": formatar_moeda(total_op2/qtd_parcelas)
                }
                
                dist_parcelas = calcular_amortizacao(qtd_parcelas)
                dados_financeiros = {
                    'acoes': acoes_fin, 
                    'total_op1': total_op1, 
                    'total_op2': total_op2,
                    'parcelas': dist_parcelas
                }
                
                with st.spinner("Construindo arquivo..."):
                    st.session_state.pptx_gerado = processar_apresentacao(template_upload, mapa_comercial, atividades_lista, "Comercial", dados_financeiros, qtd_meses_projeto)
                    st.session_state.nome_arquivo = f"Comercial_{cliente}.pptx"
                st.success("Comercial gerada com sucesso! Clique abaixo para baixar.")
                st.session_state.tentou_gerar = False

        if st.session_state.pptx_gerado and st.session_state.nome_arquivo.startswith("Comercial"):
            st.download_button("⬇️ Baixar Documento Gerado", data=st.session_state.pptx_gerado, file_name=st.session_state.nome_arquivo, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
