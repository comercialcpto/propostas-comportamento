import io
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from ferramentas.utilidades import formatar_moeda

VERDE_CPTO = RGBColor(0, 153, 116)
CINZA_ESCURO = RGBColor(64, 64, 64)


def _substituir_tags_paragrafo(paragraph, mapa):
    """
    Substitui {{TAGS}} mesmo quando o PowerPoint quebra a tag em vários runs.

    O PowerPoint frequentemente fragmenta '{{CLIENTE}}' em runs separados
    ('{{', 'CLIENTE', '}}') por causa de autocorreção/formatação. A versão
    antiga (run.text por run) falhava silenciosamente nesses casos. Aqui
    montamos o texto completo do parágrafo, substituímos, e devolvemos tudo
    no primeiro run (preservando a formatação dele).
    """
    if not paragraph.runs:
        return
    texto_completo = "".join(run.text for run in paragraph.runs)
    texto_novo = texto_completo
    for key, value in mapa.items():
        if key in texto_novo:
            texto_novo = texto_novo.replace(key, str(value))
    if texto_novo != texto_completo:
        paragraph.runs[0].text = texto_novo
        for run in paragraph.runs[1:]:
            run.text = ""


def formatar_celula_tabela(cell, texto):
    cell.text = str(texto)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = "DIN Alternate"
            run.font.size = Pt(14)


def deletar_slide(prs, slide):
    id_dict = {s.id: [i, s.rId] for i, s in enumerate(prs.slides._sldIdLst)}
    prs.part.drop_rel(id_dict[slide.slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide.slide_id][0]]


def remover_linha_tabela(table, row_idx):
    tr = table.rows[row_idx]._tr
    tr.getparent().remove(tr)


def remover_coluna_tabela(table, col_idx):
    tbl = table._tbl
    grid = tbl.tblGrid
    col = grid.gridCol_lst[col_idx]
    grid.remove(col)
    for tr in tbl.tr_lst:
        tc = tr.tc_lst[col_idx]
        tr.remove(tc)


def processar_apresentacao(template_file, mapa, atividades, tipo_doc, dados_fin=None, qtd_meses=12):
    """
    Retorna (output_BytesIO, avisos).

    'avisos' é uma lista de strings: se algo falhar no preenchimento de uma
    tabela (Gantt ou financeira), em vez de engolir o erro em silêncio,
    devolvemos um aviso para o módulo exibir via st.warning. Assim você
    descobre na hora por que uma tabela saiu vazia.
    """
    avisos = []
    prs = Presentation(template_file)
    slides_para_deletar = []

    for slide in prs.slides:
        deletar_este_slide = False

        for shape in slide.shapes:
            # Marcador de slide condicional ("Para DCS")
            if hasattr(shape, "text") and "Para DCS" in shape.text:
                if mapa.get("{{SERVICO}}", "") != "Diagnóstico (DCS/Clima/DCMA)":
                    deletar_este_slide = True
                    break

            # Substituição de tags em caixas de texto
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    _substituir_tags_paragrafo(paragraph, mapa)

            if shape.has_table:
                tbl = shape.table

                # Tags dentro de células
                for row in tbl.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            for p in cell.text_frame.paragraphs:
                                _substituir_tags_paragrafo(p, mapa)

                # Cronograma de Gantt (tabela larga, >= 12 colunas, com atividades)
                if len(tbl.columns) >= 12 and len(atividades) > 0:
                    try:
                        _preencher_gantt(prs, shape, tbl, atividades, qtd_meses)
                    except Exception as e:
                        avisos.append(f"Não foi possível montar o cronograma de Gantt: {e}")

                # Tabelas financeiras (Comercial)
                if tipo_doc == "Comercial" and dados_fin:
                    try:
                        _preencher_financeiro(tbl, dados_fin)
                    except Exception as e:
                        avisos.append(f"Não foi possível preencher uma tabela financeira: {e}")

        if deletar_este_slide:
            slides_para_deletar.append(slide)

    for slide in slides_para_deletar:
        deletar_slide(prs, slide)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output, avisos


def _preencher_gantt(prs, shape, tbl, atividades, qtd_meses):
    colunas_para_deletar = list(range(qtd_meses + 1, len(tbl.columns)))

    largura_original = shape.width
    largura_deletada = sum([tbl.columns[c].width for c in colunas_para_deletar])
    nova_largura = largura_original - largura_deletada
    shape.left = int((prs.slide_width - nova_largura) / 2)

    for c_idx in reversed(colunas_para_deletar):
        remover_coluna_tabela(tbl, c_idx)

    linhas_deletar = list(range(len(atividades) + 1, len(tbl.rows)))
    for r_idx in reversed(linhas_deletar):
        remover_linha_tabela(tbl, r_idx)

    for row_idx, atividade in enumerate(atividades):
        target_row = row_idx + 1
        if target_row < len(tbl.rows):
            row = tbl.rows[target_row]
            cell = row.cells[0]
            cell.text = atividade['nome']

            tamanho_str = len(atividade['nome'])
            fonte_tamanho = 12
            if tamanho_str > 60:
                fonte_tamanho = 8
            elif tamanho_str > 40:
                fonte_tamanho = 9
            elif tamanho_str > 20:
                fonte_tamanho = 10

            if cell.text_frame.paragraphs:
                p = cell.text_frame.paragraphs[0]
                if p.runs:
                    run = p.runs[0]
                    run.font.name = "Calibri"
                    run.font.size = Pt(fonte_tamanho)
                    run.font.color.rgb = CINZA_ESCURO

            for m_idx in range(1, len(tbl.columns)):
                if m_idx in atividade['meses']:
                    cell_mes = row.cells[m_idx]
                    cell_mes.fill.solid()
                    cell_mes.fill.fore_color.rgb = VERDE_CPTO


def _preencher_financeiro(tbl, dados_fin):
    cabecalho = tbl.rows[0].cells[0].text.strip().lower()

    if "macro" in cabecalho:
        acoes = dados_fin['acoes']
        linhas_para_deletar = []
        for idx in range(1, len(tbl.rows)):
            cell_text = tbl.rows[idx].cells[0].text.strip().lower()
            if "investimento total" in cell_text:
                formatar_celula_tabela(tbl.rows[idx].cells[1], formatar_moeda(dados_fin['total_op1']))
                formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(dados_fin['total_op2']))
            elif idx <= len(acoes):
                formatar_celula_tabela(tbl.rows[idx].cells[0], acoes[idx - 1]['nome'])
                formatar_celula_tabela(tbl.rows[idx].cells[1], formatar_moeda(acoes[idx - 1]['v1']))
                formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(acoes[idx - 1]['v2']))
            else:
                linhas_para_deletar.append(idx)
        for idx in reversed(linhas_para_deletar):
            remover_linha_tabela(tbl, idx)

    elif "meses" in cabecalho:
        parcelas = dados_fin['parcelas']
        linhas_para_deletar = []
        for idx in range(1, len(tbl.rows)):
            cell_text = tbl.rows[idx].cells[0].text.strip().lower()
            if "total" in cell_text and "investimento" not in cell_text:
                formatar_celula_tabela(tbl.rows[idx].cells[1], "100%")
                formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(dados_fin['total_op2']))
            elif idx <= len(parcelas):
                formatar_celula_tabela(tbl.rows[idx].cells[0], f"M{idx}")
                formatar_celula_tabela(tbl.rows[idx].cells[1], f"{parcelas[idx - 1]}%")
                val_calc = dados_fin['total_op2'] * (parcelas[idx - 1] / 100)
                formatar_celula_tabela(tbl.rows[idx].cells[2], formatar_moeda(val_calc))
            else:
                linhas_para_deletar.append(idx)
        for idx in reversed(linhas_para_deletar):
            remover_linha_tabela(tbl, idx)
